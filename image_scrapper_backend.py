# image_scrapper_backend.py

import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

from PIL import Image


def save_images_to_ppt(images, output_ppt):
    prs = Presentation()
    width = prs.slide_width
    height = prs.slide_height
    for img_data in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        stream = BytesIO()
        img_data.save(stream, format='PNG')
        stream.seek(0)
        slide.shapes.add_picture(stream, 0, 0, width=width, height=height)
    prs.save(output_ppt)


def extract_images_by_keyword(pdf_path, keyword, output_ppt, progress_data=None):
    images = []
    skipped = 0
    logs = []

    with pdfplumber.open(pdf_path) as pdf:
        total = len(pdf.pages)
        for idx, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            if keyword.lower() in text.lower() and len(page.images) > 1:
                page_image = page.to_image(resolution=100).original
                images.append(page_image)
                logs.append(f"🔍 Match + image found on page {idx + 1}")
            else:
                skipped += 1
            if progress_data is not None:
                progress_data["progress"] = int(((idx + 1) / total) * 100)

    if images:
        save_images_to_ppt(images, output_ppt)
        logs.append(f"✅ Saved {len(images)} valid slides to: {output_ppt}")
    logs.append(f"⚠️ Skipped {skipped} pages.")
    return logs


def extract_images_by_page(pdf_path, page_nums, output_ppt, progress_data=None):
    images = []
    skipped = 0
    logs = []

    try:
        pages = [int(p.strip()) for p in page_nums.split(',') if p.strip().isdigit()]
    except Exception as e:
        return [f"❌ Error parsing page numbers: {e}"]

    with pdfplumber.open(pdf_path) as pdf:
        total = len(pages)
        for idx, page_num in enumerate(pages):
            if 0 < page_num <= len(pdf.pages):
                page = pdf.pages[page_num - 1]
                if len(page.images) > 1:
                    page_image = page.to_image(resolution=100).original
                    images.append(page_image)
                    logs.append(f"📄 Snapshot taken from page {page_num}")
                else:
                    skipped += 1
            if progress_data is not None:
                progress_data["progress"] = int(((idx + 1) / total) * 100)

    if images:
        save_images_to_ppt(images, output_ppt)
        logs.append(f"✅ Saved {len(images)} valid slides to: {output_ppt}")
    logs.append(f"⚠️ Skipped {skipped} pages.")
    return logs
